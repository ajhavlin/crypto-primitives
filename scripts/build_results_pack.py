#!/usr/bin/env python3
"""Build a 7-slide results replacement pack (PPTX) from benchmark figure PNGs.

Usage:
    python3 scripts/build_results_pack.py \
        --fig-dir figures/presentation \
        --out /mnt/c/Users/ajhav/Downloads/revised_batch_proofs_strategy_results_pack.pptx \
        --title Results
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

EMU_PER_INCH = 914400

# 16:9 widescreen
SLIDE_W = Emu(int(13.333 * EMU_PER_INCH))
SLIDE_H = Emu(int(7.5 * EMU_PER_INCH))

TOKENS = {
    "accent": RGBColor(48, 93, 255),
    "bg_dark": RGBColor(5, 14, 42),
    "bg": RGBColor(248, 250, 252),
    "text_dark": RGBColor(18, 25, 38),
    "text_muted": RGBColor(88, 96, 112),
    "card": RGBColor(241, 244, 249),
    "card_border": RGBColor(219, 225, 236),
    "sticky": RGBColor(255, 247, 214),
    "sticky_border": RGBColor(234, 222, 172),
    "progress_bg": RGBColor(234, 240, 255),
}


@dataclass(frozen=True)
class FigureSlide:
    filename: str
    title: str
    subtitle: str
    callout_header: str
    callout_body: str
    progress: str


FIGURE_SLIDES: Tuple[FigureSlide, ...] = (
    FigureSlide(
        filename="proof_size_vs_k_clustered.png",
        title="Proof Size - Clustered",
        subtitle="One graph per slide. Shared y-scale with random for clean comparison.",
        callout_header="🎯 Clustered wins clearly",
        callout_body="Coset pruning removes duplicated siblings, giving the largest byte savings at high k.",
        progress="Clustered 1/3",
    ),
    FigureSlide(
        filename="prover_time_vs_k_clustered.png",
        title="Prover Time - Clustered",
        subtitle="Same visual frame to support narrative pacing.",
        callout_header="⏱️ Prover overhead is controlled",
        callout_body="Extraction/reconstruction costs stay modest versus structural proof-size gains.",
        progress="Clustered 2/3",
    ),
    FigureSlide(
        filename="verifier_time_vs_k_clustered.png",
        title="Verifier Time - Clustered",
        subtitle="Consistent scale and typography to reduce audience load.",
        callout_header="⚖️ Hashing still dominates",
        callout_body="Verifier trend remains hash-heavy; metadata optimizations have a smaller effect.",
        progress="Clustered 3/3",
    ),
    FigureSlide(
        filename="proof_size_vs_k_random.png",
        title="Proof Size - Random",
        subtitle="Now random queries under the same axis policy.",
        callout_header="🧭 Savings persist in random",
        callout_body="Overlap is lower than clustered, but coset still reduces proof bytes materially.",
        progress="Random 1/3",
    ),
    FigureSlide(
        filename="prover_time_vs_k_random.png",
        title="Prover Time - Random",
        subtitle="Directly comparable to clustered due to shared y-ranges.",
        callout_header="🔎 Runtime tradeoff is visible",
        callout_body="Prover-time gap remains moderate relative to proof-size improvements.",
        progress="Random 2/3",
    ),
    FigureSlide(
        filename="verifier_time_vs_k_random.png",
        title="Verifier Time - Random",
        subtitle="Final evidence slide in the same visual grammar.",
        callout_header="✅ End-to-end story closes",
        callout_body="Verifier cost tracks hashing work; optimization impact is stable across query styles.",
        progress="Random 3/3",
    ),
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--fig-dir",
        type=Path,
        default=Path("figures/presentation"),
        help="Directory containing required graph PNGs.",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=Path("/mnt/c/Users/ajhav/Downloads/revised_batch_proofs_strategy_results_pack.pptx"),
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--title",
        type=str,
        default="Results",
        help="Section divider title word.",
    )
    parser.add_argument(
        "--readme",
        type=Path,
        default=Path("/mnt/c/Users/ajhav/Downloads/revised_batch_proofs_strategy_results_pack_README.txt"),
        help="Companion README/TXT path.",
    )
    return parser.parse_args()


def ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def preflight_figure_assets(fig_dir: Path, names: Sequence[str]) -> List[Path]:
    resolved: List[Path] = []
    missing: List[str] = []
    sizes: List[Tuple[int, int]] = []

    for name in names:
        candidate = fig_dir / name
        if not candidate.exists():
            missing.append(name)
            continue
        resolved.append(candidate)
        with Image.open(candidate) as im:
            sizes.append((im.width, im.height))

    if missing:
        raise FileNotFoundError(f"Missing figure files in {fig_dir}: {', '.join(missing)}")

    heights = {h for _, h in sizes}
    if len(heights) > 1:
        raise ValueError(f"Figure heights are inconsistent: {sorted(heights)}")

    widths = [w for w, _ in sizes]
    w_min = min(widths)
    w_max = max(widths)
    if w_max > 1.35 * w_min:
        raise ValueError(
            "Figure widths are too inconsistent for a single visual family "
            f"(min={w_min}, max={w_max}, ratio={w_max / w_min:.2f})."
        )

    return resolved


def add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_rounded_box(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    fill: RGBColor,
    line: RGBColor | None = None,
    radius_adjust: float = 0.12,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    # Smaller corner radius for a sleek card look
    shape.adjustments[0] = radius_adjust
    return shape


def place_picture_contain(slide, image_path: Path, left: Emu, top: Emu, width: Emu, height: Emu) -> None:
    with Image.open(image_path) as im:
        img_w, img_h = im.size
    scale = min(float(width) / img_w, float(height) / img_h)
    pic_w = Emu(int(img_w * scale))
    pic_h = Emu(int(img_h * scale))
    pic_left = Emu(int(left + (width - pic_w) / 2))
    pic_top = Emu(int(top + (height - pic_h) / 2))
    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def build_divider_slide(prs: Presentation, title_word: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TOKENS["bg_dark"]
    bg.line.fill.background()

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0),
        Emu(int(5.6 * EMU_PER_INCH)),
        SLIDE_W,
        Emu(int(0.16 * EMU_PER_INCH)),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = TOKENS["accent"]
    band.line.fill.background()

    add_textbox(
        slide,
        Emu(int(1.1 * EMU_PER_INCH)),
        Emu(int(2.4 * EMU_PER_INCH)),
        Emu(int(11.2 * EMU_PER_INCH)),
        Emu(int(1.6 * EMU_PER_INCH)),
        title_word,
        size=66,
        bold=True,
        color=RGBColor(248, 250, 255),
    )

    add_textbox(
        slide,
        Emu(int(1.1 * EMU_PER_INCH)),
        Emu(int(4.15 * EMU_PER_INCH)),
        Emu(int(8.2 * EMU_PER_INCH)),
        Emu(int(0.6 * EMU_PER_INCH)),
        "Measured with one graph per slide for cleaner storytelling.",
        size=20,
        color=RGBColor(193, 203, 226),
    )

    chip = add_rounded_box(
        slide,
        Emu(int(12.35 * EMU_PER_INCH)),
        Emu(int(0.35 * EMU_PER_INCH)),
        Emu(int(0.28 * EMU_PER_INCH)),
        Emu(int(0.28 * EMU_PER_INCH)),
        fill=TOKENS["accent"],
        line=None,
        radius_adjust=0.35,
    )
    chip.line.fill.background()


def build_metric_slide(prs: Presentation, spec: FigureSlide, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TOKENS["bg"]
    bg.line.fill.background()

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(int(0.08 * EMU_PER_INCH)),
        Emu(int(0.5 * EMU_PER_INCH)),
        Emu(int(0.05 * EMU_PER_INCH)),
        Emu(int(6.5 * EMU_PER_INCH)),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TOKENS["accent"]
    accent_bar.line.fill.background()

    add_textbox(
        slide,
        Emu(int(0.5 * EMU_PER_INCH)),
        Emu(int(0.35 * EMU_PER_INCH)),
        Emu(int(8.6 * EMU_PER_INCH)),
        Emu(int(0.72 * EMU_PER_INCH)),
        spec.title,
        size=34,
        bold=True,
        color=TOKENS["text_dark"],
    )
    add_textbox(
        slide,
        Emu(int(0.5 * EMU_PER_INCH)),
        Emu(int(0.96 * EMU_PER_INCH)),
        Emu(int(8.9 * EMU_PER_INCH)),
        Emu(int(0.52 * EMU_PER_INCH)),
        spec.subtitle,
        size=16,
        color=TOKENS["text_muted"],
    )

    card_left = Emu(int(0.48 * EMU_PER_INCH))
    card_top = Emu(int(1.42 * EMU_PER_INCH))
    card_w = Emu(int(9.65 * EMU_PER_INCH))
    card_h = Emu(int(5.55 * EMU_PER_INCH))

    add_rounded_box(
        slide,
        card_left,
        card_top,
        card_w,
        card_h,
        fill=TOKENS["card"],
        line=TOKENS["card_border"],
        radius_adjust=0.06,
    )

    pic_margin = Emu(int(0.18 * EMU_PER_INCH))
    place_picture_contain(
        slide,
        image_path,
        Emu(int(card_left + pic_margin)),
        Emu(int(card_top + pic_margin)),
        Emu(int(card_w - 2 * pic_margin)),
        Emu(int(card_h - 2 * pic_margin)),
    )

    sticky_left = Emu(int(10.45 * EMU_PER_INCH))
    sticky_top = Emu(int(1.65 * EMU_PER_INCH))
    sticky_w = Emu(int(2.62 * EMU_PER_INCH))
    sticky_h = Emu(int(2.76 * EMU_PER_INCH))

    add_rounded_box(
        slide,
        sticky_left,
        sticky_top,
        sticky_w,
        sticky_h,
        fill=TOKENS["sticky"],
        line=TOKENS["sticky_border"],
        radius_adjust=0.09,
    )

    add_textbox(
        slide,
        Emu(int(sticky_left + 0.18 * EMU_PER_INCH)),
        Emu(int(sticky_top + 0.16 * EMU_PER_INCH)),
        Emu(int(sticky_w - 0.35 * EMU_PER_INCH)),
        Emu(int(0.7 * EMU_PER_INCH)),
        spec.callout_header,
        size=16,
        bold=True,
        color=TOKENS["text_dark"],
    )
    add_textbox(
        slide,
        Emu(int(sticky_left + 0.18 * EMU_PER_INCH)),
        Emu(int(sticky_top + 0.74 * EMU_PER_INCH)),
        Emu(int(sticky_w - 0.35 * EMU_PER_INCH)),
        Emu(int(1.9 * EMU_PER_INCH)),
        spec.callout_body,
        size=14,
        color=RGBColor(61, 67, 79),
    )

    tag = add_rounded_box(
        slide,
        Emu(int(10.52 * EMU_PER_INCH)),
        Emu(int(4.88 * EMU_PER_INCH)),
        Emu(int(2.5 * EMU_PER_INCH)),
        Emu(int(0.54 * EMU_PER_INCH)),
        fill=TOKENS["progress_bg"],
        line=TOKENS["accent"],
        radius_adjust=0.22,
    )
    tag.line.width = Pt(1.1)
    add_textbox(
        slide,
        Emu(int(10.62 * EMU_PER_INCH)),
        Emu(int(5.01 * EMU_PER_INCH)),
        Emu(int(2.3 * EMU_PER_INCH)),
        Emu(int(0.34 * EMU_PER_INCH)),
        spec.progress,
        size=12,
        bold=True,
        color=TOKENS["accent"],
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Emu(int(10.45 * EMU_PER_INCH)),
        Emu(int(6.43 * EMU_PER_INCH)),
        Emu(int(2.65 * EMU_PER_INCH)),
        Emu(int(0.4 * EMU_PER_INCH)),
        "Results chapter redesign",
        size=11,
        color=RGBColor(114, 121, 136),
        align=PP_ALIGN.RIGHT,
    )


def build_readme(readme_path: Path, pptx_path: Path, fig_dir: Path) -> None:
    text = f"""Results Replacement Pack - Quick Insert Guide

Output deck:
{pptx_path}

Source figures:
{fig_dir}

Insertion steps:
1) Open the original Keynote deck and this PPTX side-by-side.
2) In the original deck, remove the existing Results block (old divider + old metric slides).
3) Import all 7 slides from the PPTX in order:
   - Results divider
   - Proof Size - Clustered
   - Prover Time - Clustered
   - Verifier Time - Clustered
   - Proof Size - Random
   - Prover Time - Random
   - Verifier Time - Random
4) Keep downstream slides (Migration/Next/Questions) as-is.

Style conformance checklist:
- Section architecture: full-bleed divider with a single large chapter word.
- Visual system: one accent color, neutral palette, fixed typography hierarchy, repeated rounded cards.
- Controlled density: one graph per slide with generous whitespace.
- Progressive disclosure: clustered trio first, then random trio in an identical frame.
- Embedded annotations: one modern sticky-note insight callout per metric slide.
- Tasteful informality: subtle semantic emoji in callout headers.

Validation:
The redesigned chapter matches the requested benchmark-slide style principles and keeps comparisons clean via one-graph-per-slide pacing. Next step: paste this 7-slide block into the main Keynote deck.
"""
    readme_path.write_text(text, encoding="utf-8")


def build_pack(fig_dir: Path, out_path: Path, title: str, readme_path: Path) -> None:
    required = [s.filename for s in FIGURE_SLIDES]
    resolved = preflight_figure_assets(fig_dir, required)
    figure_map = {p.name: p for p in resolved}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_divider_slide(prs, title)
    for spec in FIGURE_SLIDES:
        build_metric_slide(prs, spec, figure_map[spec.filename])

    ensure_parent(out_path)
    prs.save(out_path)

    ensure_parent(readme_path)
    build_readme(readme_path, out_path, fig_dir)


def main() -> None:
    args = parse_args()
    build_pack(args.fig_dir, args.out, args.title, args.readme)
    print(f"Wrote {args.out}")
    print(f"Wrote {args.readme}")


if __name__ == "__main__":
    main()
